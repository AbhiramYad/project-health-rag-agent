"""
Monthly Portfolio Synthesis & Executive Presentation Generator
===============================================================
Reads all weekly health report JSON files from the outputs/ directory,
synthesizes cross-project trends using Gemini LLM, and generates a
5-7 slide executive PowerPoint presentation (.pptx).

Usage:
    python synthesis.py

Requirements:
    pip install python-pptx google-genai python-dotenv
    Add your key to .env file: GEMINI_API_KEY=your_key_here
"""

import os
import sys
import json
import glob
from datetime import datetime
from collections import defaultdict
from dotenv import load_dotenv
from google import genai
from google.genai import types

# python-pptx imports
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

load_dotenv()

if sys.stdout.encoding and sys.stdout.encoding.lower() != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8')

# ─────────────────────────────────────────────
# CONFIG
# ─────────────────────────────────────────────
GEMINI_MODEL  = "gemini-2.5-flash"
OUTPUT_DIR    = "outputs"
PPTX_OUTPUT   = os.path.join(OUTPUT_DIR, "monthly_presentation.pptx")

# Brand colours
C_DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
C_ACCENT    = RGBColor(0x16, 0x21, 0x3E)   # dark blue panel
C_HIGHLIGHT = RGBColor(0x0F, 0x3E, 0x80)   # vivid blue accent
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xE8, 0xF4, 0xFF)
C_GREEN     = RGBColor(0x27, 0xAE, 0x60)
C_AMBER     = RGBColor(0xF3, 0x9C, 0x12)
C_RED       = RGBColor(0xC0, 0x39, 0x2B)
C_GRAY      = RGBColor(0x95, 0xA5, 0xA6)

RAG_COLORS  = {"Green": C_GREEN, "Amber": C_AMBER, "Red": C_RED, "Unknown": C_GRAY}
RAG_EMOJI   = {"Green": "GREEN", "Amber": "AMBER", "Red": "RED", "Unknown": "UNKNOWN"}

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────
# DATA LOADING
# ─────────────────────────────────────────────

def load_all_reports() -> list[dict]:
    """Load all weekly JSON reports from the outputs/ directory."""
    pattern = os.path.join(OUTPUT_DIR, "*", "*_health_report.json")
    files   = sorted(glob.glob(pattern))
    reports = []
    for f in files:
        try:
            with open(f, encoding="utf-8") as fh:
                data = json.load(fh)
                data["_source_file"] = f
                reports.append(data)
        except Exception as e:
            print(f"[!] Could not load {f}: {e}")
    return reports


def aggregate_reports(reports: list[dict]) -> dict:
    """Aggregate statistics across all reports for synthesis."""
    by_project  = defaultdict(list)
    by_date     = defaultdict(list)
    all_risks   = []
    all_actions = []

    for r in reports:
        by_project[r["project_key"]].append(r)
        by_date[r["report_date"]].append(r)
        all_risks   += r.get("key_risks", [])
        all_actions += r.get("recommended_actions", [])

    # Latest report per project
    latest = {}
    for key, reps in by_project.items():
        latest[key] = sorted(reps, key=lambda x: x["report_date"])[-1]

    # RAG trend per project (chronological)
    trends = {}
    for key, reps in by_project.items():
        trends[key] = [
            {"date": r["report_date"], "rag": r["rag_status"], "score": r["final_score"]}
            for r in sorted(reps, key=lambda x: x["report_date"])
        ]

    # Portfolio-level stats
    total_projects = len(latest)
    rag_counts = {"Red": 0, "Amber": 0, "Green": 0}
    for r in latest.values():
        rag_counts[r["rag_status"]] = rag_counts.get(r["rag_status"], 0) + 1

    dates = sorted(by_date.keys())

    return {
        "reports":        reports,
        "by_project":     dict(by_project),
        "latest":         latest,
        "trends":         trends,
        "all_risks":      all_risks,
        "all_actions":    all_actions,
        "total_projects": total_projects,
        "rag_counts":     rag_counts,
        "report_dates":   dates,
        "month_label":    datetime.now().strftime("%B %Y"),
    }


# ─────────────────────────────────────────────
# LLM SYNTHESIS
# ─────────────────────────────────────────────

def build_synthesis_prompt(agg: dict) -> str:
    """Build a prompt asking Gemini to synthesize cross-project insights."""
    project_summaries = []
    for key, r in agg["latest"].items():
        project_summaries.append(
            f"- {key}: RAG={r['rag_status']} | {r['pct_complete']:.0f}% complete | "
            f"Score={r['final_score']:.2f}/2.0 | At Risk={r['at_risk']} | "
            f"Stage={r['project_stage']}\n"
            f"  Summary: {r.get('summary','N/A')}"
        )

    risks_text    = "\n".join(f"- {r}" for r in agg["all_risks"][:15])
    actions_text  = "\n".join(f"- {a}" for a in agg["all_actions"][:10])

    return f"""You are a senior delivery director at Zycus Professional Services preparing a monthly executive briefing.

PORTFOLIO OVERVIEW - {agg['month_label']}
==========================================
Total Projects Monitored : {agg['total_projects']}
Red    : {agg['rag_counts'].get('Red', 0)} project(s)
Amber  : {agg['rag_counts'].get('Amber', 0)} project(s)
Green  : {agg['rag_counts'].get('Green', 0)} project(s)
Report Dates Covered: {', '.join(agg['report_dates'])}

LATEST PROJECT STATUS
======================
{chr(10).join(project_summaries)}

ALL IDENTIFIED RISKS (across all reports)
==========================================
{risks_text}

ALL RECOMMENDED ACTIONS (across all reports)
=============================================
{actions_text}

INSTRUCTIONS:
Generate a concise monthly executive synthesis. Output EXACTLY in this format:

PORTFOLIO_HEADLINE: [1 sentence capturing the overall portfolio health this month]

TREND_1: [Cross-project trend observation 1 - must reference specific projects]
TREND_2: [Cross-project trend observation 2]
TREND_3: [Cross-project trend observation 3]

EMERGING_RISK_1: [Systemic risk 1 that affects multiple projects or has escalating pattern]
EMERGING_RISK_2: [Systemic risk 2]
EMERGING_RISK_3: [Systemic risk 3]

RECOMMENDATION_1: [Strategic recommendation for leadership - actionable, specific]
RECOMMENDATION_2: [Strategic recommendation 2]
RECOMMENDATION_3: [Strategic recommendation 3]
RECOMMENDATION_4: [Strategic recommendation 4]
RECOMMENDATION_5: [Strategic recommendation 5]

CONCLUSION: [2-3 sentence closing statement a VP would say to a client]
"""


def call_gemini(prompt: str, api_key: str) -> str:
    client = genai.Client(api_key=api_key)
    response = client.models.generate_content(
        model=GEMINI_MODEL,
        contents=prompt,
        config=types.GenerateContentConfig(temperature=0.3, max_output_tokens=2048),
    )
    return response.text.strip()


def parse_synthesis(text: str) -> dict:
    """Parse the structured synthesis LLM response."""
    import re
    result = {
        "headline": "", "trends": [], "emerging_risks": [],
        "recommendations": [], "conclusion": "", "raw": text
    }

    def extract(key):
        m = re.search(rf"{key}:\s*(.+?)(?=\n[A-Z_]+:|$)", text, re.DOTALL)
        return m.group(1).strip() if m else ""

    result["headline"]    = extract("PORTFOLIO_HEADLINE")
    result["conclusion"]  = extract("CONCLUSION")

    for i in range(1, 5):
        v = extract(f"TREND_{i}")
        if v: result["trends"].append(v)

    for i in range(1, 5):
        v = extract(f"EMERGING_RISK_{i}")
        if v: result["emerging_risks"].append(v)

    for i in range(1, 7):
        v = extract(f"RECOMMENDATION_{i}")
        if v: result["recommendations"].append(v)

    return result


# ─────────────────────────────────────────────
# PPTX HELPERS
# ─────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=None,
                align=PP_ALIGN.LEFT, italic=False):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color or C_WHITE
    return txBox


def add_rect(slide, left, top, width, height, fill_color: RGBColor, transparency=0):
    """Add a filled rectangle (panel) to a slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def rag_badge_text(rag: str) -> str:
    icons = {"Red": "[RED]", "Amber": "[AMBER]", "Green": "[GREEN]", "Unknown": "[?]"}
    return icons.get(rag, "[?]")


# ─────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────

def build_cover_slide(prs, agg: dict, synthesis: dict):
    """Slide 1: Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, C_DARK_BG)

    # Top accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_HIGHLIGHT)

    # Zycus logo text
    add_textbox(slide, "ZYCUS", Inches(0.5), Inches(0.3), Inches(3), Inches(0.5),
                font_size=14, bold=True, color=C_HIGHLIGHT)

    # Main title
    add_textbox(slide, "Project Portfolio\nHealth Report",
                Inches(0.5), Inches(1.5), Inches(8), Inches(2),
                font_size=44, bold=True, color=C_WHITE)

    # Subtitle
    add_textbox(slide, agg["month_label"],
                Inches(0.5), Inches(3.6), Inches(6), Inches(0.8),
                font_size=26, bold=False, color=C_HIGHLIGHT)

    # Headline
    if synthesis["headline"]:
        add_textbox(slide, synthesis["headline"],
                    Inches(0.5), Inches(4.5), Inches(10), Inches(1.2),
                    font_size=15, italic=True, color=C_LIGHT)

    # Stats row
    stats = [
        (str(agg["rag_counts"].get("Red", 0)),   "RED",   C_RED,   Inches(9.5)),
        (str(agg["rag_counts"].get("Amber", 0)), "AMBER", C_AMBER, Inches(10.7)),
        (str(agg["rag_counts"].get("Green", 0)), "GREEN", C_GREEN, Inches(11.9)),
    ]
    for val, label, color, left in stats:
        add_rect(slide, left, Inches(5.8), Inches(1.0), Inches(1.2), color)
        add_textbox(slide, val, left, Inches(5.85), Inches(1.0), Inches(0.6),
                    font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, left, Inches(6.5), Inches(1.0), Inches(0.4),
                    font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, f"Prepared by: AI Project Health Agent  |  Confidential",
                Inches(0.5), Inches(7.0), Inches(9), Inches(0.4),
                font_size=9, color=C_GRAY)

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), C_HIGHLIGHT)


def build_executive_summary_slide(prs, agg: dict, synthesis: dict):
    """Slide 2: Executive Summary — RAG Dashboard"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_HIGHLIGHT)

    add_textbox(slide, "Executive Summary", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
                font_size=22, bold=True, color=C_WHITE)
    add_textbox(slide, "Portfolio RAG Status Dashboard",
                Inches(0.5), Inches(0.8), Inches(10), Inches(0.4),
                font_size=13, color=C_GRAY)

    # Project cards
    projects = list(agg["latest"].values())
    card_width  = Inches(5.8)
    card_height = Inches(2.5)
    positions   = [
        (Inches(0.4), Inches(1.4)),
        (Inches(6.4), Inches(1.4)),
        (Inches(0.4), Inches(4.1)),
        (Inches(6.4), Inches(4.1)),
    ]

    for i, r in enumerate(projects[:4]):
        if i >= len(positions):
            break
        lft, tp = positions[i]
        rag      = r["rag_status"]
        bg_color = C_ACCENT

        add_rect(slide, lft, tp, card_width, card_height, bg_color)
        # RAG indicator stripe
        add_rect(slide, lft, tp, Inches(0.18), card_height, RAG_COLORS.get(rag, C_GRAY))

        name = r["project_key"].replace("_", " ")
        add_textbox(slide, name, lft + Inches(0.3), tp + Inches(0.12),
                    Inches(3.5), Inches(0.45), font_size=14, bold=True, color=C_WHITE)

        add_textbox(slide, rag_badge_text(rag),
                    lft + Inches(4.2), tp + Inches(0.12), Inches(1.4), Inches(0.45),
                    font_size=13, bold=True, color=RAG_COLORS.get(rag, C_GRAY),
                    align=PP_ALIGN.RIGHT)

        details = (
            f"PM: {r.get('project_manager','N/A')}  |  Stage: {r.get('project_stage','N/A')}\n"
            f"Progress: {r.get('pct_complete',0):.0f}%  |  Score: {r.get('final_score',0):.2f}/2.0\n"
            f"{r.get('summary','')[:120]}..."
        )
        add_textbox(slide, details, lft + Inches(0.3), tp + Inches(0.65),
                    Inches(5.3), Inches(1.7), font_size=10, color=C_LIGHT)

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), C_HIGHLIGHT)


def build_trends_slide(prs, agg: dict, synthesis: dict):
    """Slide 3: Cross-Project Trends"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_HIGHLIGHT)

    add_textbox(slide, "Schedule & Delivery Trends", Inches(0.5), Inches(0.2),
                Inches(10), Inches(0.6), font_size=22, bold=True, color=C_WHITE)
    add_textbox(slide, "Patterns identified across all monitored projects this month",
                Inches(0.5), Inches(0.8), Inches(10), Inches(0.4), font_size=13, color=C_GRAY)

    # Trend blocks
    top = Inches(1.4)
    for i, trend in enumerate(synthesis["trends"][:4]):
        add_rect(slide, Inches(0.4), top, Inches(12.4), Inches(1.0), C_ACCENT)
        add_rect(slide, Inches(0.4), top, Inches(0.18), Inches(1.0), C_HIGHLIGHT)
        num_label = f"0{i+1}"
        add_textbox(slide, num_label, Inches(0.7), top + Inches(0.12),
                    Inches(0.8), Inches(0.75), font_size=24, bold=True, color=C_HIGHLIGHT)
        add_textbox(slide, trend, Inches(1.6), top + Inches(0.12),
                    Inches(11.0), Inches(0.75), font_size=12, color=C_WHITE)
        top += Inches(1.15)

    # RAG score comparison table
    add_textbox(slide, "Project Score Comparison", Inches(0.5), top + Inches(0.1),
                Inches(6), Inches(0.4), font_size=13, bold=True, color=C_LIGHT)
    top += Inches(0.55)
    for key, r in agg["latest"].items():
        bar_w = Inches(6.0 * (r["final_score"] / 2.0))
        add_rect(slide, Inches(0.5), top, bar_w, Inches(0.3), RAG_COLORS.get(r["rag_status"], C_GRAY))
        label = f"{key.replace('_',' ')}  {r['final_score']:.2f}/2.0"
        add_textbox(slide, label, Inches(0.6), top, Inches(7), Inches(0.3),
                    font_size=10, color=C_WHITE)
        top += Inches(0.42)

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), C_HIGHLIGHT)


def build_risks_slide(prs, agg: dict, synthesis: dict):
    """Slide 4: Emerging Risks & Blockers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_RED)

    add_textbox(slide, "Emerging Risks & Blockers", Inches(0.5), Inches(0.2),
                Inches(10), Inches(0.6), font_size=22, bold=True, color=C_WHITE)
    add_textbox(slide, "Systemic risks requiring immediate leadership attention",
                Inches(0.5), Inches(0.8), Inches(10), Inches(0.4), font_size=13, color=C_GRAY)

    top = Inches(1.4)
    for risk in synthesis["emerging_risks"][:5]:
        add_rect(slide, Inches(0.4), top, Inches(12.4), Inches(0.9), C_ACCENT)
        add_rect(slide, Inches(0.4), top, Inches(0.18), Inches(0.9), C_RED)
        add_textbox(slide, "!", Inches(0.65), top + Inches(0.08),
                    Inches(0.4), Inches(0.7), font_size=20, bold=True, color=C_RED,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, risk, Inches(1.2), top + Inches(0.1),
                    Inches(11.4), Inches(0.7), font_size=12, color=C_WHITE)
        top += Inches(1.05)

    # At Risk summary box
    add_rect(slide, Inches(0.4), top + Inches(0.1), Inches(12.4), Inches(0.75), C_ACCENT)
    at_risk_projects = [k for k, r in agg["latest"].items() if r.get("at_risk","").lower() == "high"]
    summary = f"Projects flagged At Risk = HIGH: {', '.join(at_risk_projects) or 'None'}"
    add_textbox(slide, summary, Inches(0.7), top + Inches(0.18), Inches(12.0), Inches(0.5),
                font_size=13, bold=True, color=C_AMBER)

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), C_RED)


def build_deepdive_slide(prs, agg: dict):
    """Slide 5: Project Deep-Dives"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_HIGHLIGHT)

    add_textbox(slide, "Project Deep-Dives", Inches(0.5), Inches(0.2),
                Inches(10), Inches(0.6), font_size=22, bold=True, color=C_WHITE)
    add_textbox(slide, "Latest status snapshot per project",
                Inches(0.5), Inches(0.8), Inches(10), Inches(0.4), font_size=13, color=C_GRAY)

    col_w   = Inches(6.2)
    col_gap = Inches(0.5)
    top_start = Inches(1.4)

    for idx, (key, r) in enumerate(agg["latest"].items()):
        col   = idx % 2
        row   = idx // 2
        left  = Inches(0.4) + col * (col_w + col_gap)
        top   = top_start + row * Inches(3.1)
        rag   = r["rag_status"]

        add_rect(slide, left, top, col_w, Inches(2.8), C_ACCENT)
        add_rect(slide, left, top, col_w, Inches(0.4), RAG_COLORS.get(rag, C_GRAY))

        name = key.replace("_", " ")
        add_textbox(slide, f"{name}  [{rag}]", left + Inches(0.15), top + Inches(0.05),
                    Inches(5.8), Inches(0.35), font_size=13, bold=True, color=C_WHITE)

        info_lines = [
            f"PM: {r.get('project_manager','N/A')}",
            f"Stage: {r.get('project_stage','N/A')}",
            f"Progress: {r.get('pct_complete',0):.0f}% (expected ~{r.get('expected_pct',0):.0f}%)",
            f"Tasks: {r.get('task_counts',{}).get('completed',0)} done / {r.get('task_counts',{}).get('total',0)} total",
        ]
        if r.get("key_risks"):
            info_lines.append(f"Top Risk: {r['key_risks'][0][:65]}...")

        add_textbox(slide, "\n".join(info_lines), left + Inches(0.15), top + Inches(0.5),
                    Inches(5.8), Inches(2.1), font_size=10, color=C_LIGHT)

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), C_HIGHLIGHT)


def build_recommendations_slide(prs, synthesis: dict):
    """Slide 6: Strategic Recommendations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_GREEN)

    add_textbox(slide, "Strategic Recommendations", Inches(0.5), Inches(0.2),
                Inches(10), Inches(0.6), font_size=22, bold=True, color=C_WHITE)
    add_textbox(slide, "Actions recommended for leadership this month",
                Inches(0.5), Inches(0.8), Inches(10), Inches(0.4), font_size=13, color=C_GRAY)

    top = Inches(1.4)
    for i, rec in enumerate(synthesis["recommendations"][:5]):
        add_rect(slide, Inches(0.4), top, Inches(12.4), Inches(0.9), C_ACCENT)
        add_rect(slide, Inches(0.4), top, Inches(0.18), Inches(0.9), C_GREEN)
        add_textbox(slide, str(i + 1), Inches(0.65), top + Inches(0.08),
                    Inches(0.4), Inches(0.7), font_size=20, bold=True, color=C_GREEN,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, rec, Inches(1.2), top + Inches(0.1),
                    Inches(11.4), Inches(0.7), font_size=12, color=C_WHITE)
        top += Inches(1.05)

    # Closing statement
    if synthesis["conclusion"]:
        add_rect(slide, Inches(0.4), top + Inches(0.1), Inches(12.4), Inches(0.9), C_HIGHLIGHT)
        add_textbox(slide, synthesis["conclusion"],
                    Inches(0.7), top + Inches(0.18), Inches(12.0), Inches(0.65),
                    font_size=12, italic=True, color=C_LIGHT)

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), C_GREEN)


def build_appendix_slide(prs, agg: dict):
    """Slide 7: Appendix — Full Weekly RAG Log"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_HIGHLIGHT)

    add_textbox(slide, "Appendix: Weekly RAG Log", Inches(0.5), Inches(0.2),
                Inches(10), Inches(0.6), font_size=22, bold=True, color=C_WHITE)
    add_textbox(slide, "Complete record of all weekly health assessments",
                Inches(0.5), Inches(0.8), Inches(10), Inches(0.4), font_size=13, color=C_GRAY)

    # Table header
    col_positions = [Inches(0.5), Inches(3.2), Inches(5.5), Inches(7.3), Inches(9.5), Inches(11.2)]
    headers = ["Date", "Project", "RAG", "Score", "% Done", "At Risk"]
    add_rect(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.4), C_HIGHLIGHT)
    for h, lft in zip(headers, col_positions):
        add_textbox(slide, h, lft, Inches(1.32), Inches(2.5), Inches(0.36),
                    font_size=10, bold=True, color=C_WHITE)

    top = Inches(1.8)
    alt = False
    for r in sorted(agg["reports"], key=lambda x: (x["report_date"], x["project_key"])):
        bg = C_ACCENT if alt else C_DARK_BG
        add_rect(slide, Inches(0.4), top, Inches(12.5), Inches(0.37), bg)
        row_vals = [
            r["report_date"],
            r["project_key"].replace("_", " ")[:18],
            r["rag_status"],
            f"{r['final_score']:.2f}",
            f"{r.get('pct_complete',0):.0f}%",
            r.get("at_risk", "N/A"),
        ]
        rag_color = RAG_COLORS.get(r["rag_status"], C_WHITE)
        for val, lft, h in zip(row_vals, col_positions, headers):
            color = rag_color if h == "RAG" else C_LIGHT
            add_textbox(slide, val, lft, top + Inches(0.02), Inches(2.5), Inches(0.33),
                        font_size=9, color=color)
        top += Inches(0.4)
        alt = not alt
        if top > Inches(7.0):
            break

    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), C_HIGHLIGHT)


# ─────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────

def main():
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        print("ERROR: GEMINI_API_KEY not set. Add it to your .env file.")
        return

    print("=" * 60)
    print("  Monthly Portfolio Synthesis")
    print("=" * 60)

    # 1. Load all reports
    print("\n[1] Loading all weekly reports...")
    reports = load_all_reports()
    if not reports:
        print("ERROR: No report JSON files found in outputs/. Run agent.py first.")
        return
    print(f"    Loaded {len(reports)} report(s) across {len(set(r['project_key'] for r in reports))} project(s).")

    # 2. Aggregate
    print("[2] Aggregating cross-project statistics...")
    agg = aggregate_reports(reports)
    print(f"    Portfolio: {agg['rag_counts'].get('Red',0)} Red | "
          f"{agg['rag_counts'].get('Amber',0)} Amber | "
          f"{agg['rag_counts'].get('Green',0)} Green")

    # 3. LLM synthesis
    print("[3] Calling Gemini for executive synthesis...")
    prompt    = build_synthesis_prompt(agg)
    llm_text  = call_gemini(prompt, api_key)
    synthesis = parse_synthesis(llm_text)
    print(f"    Headline: {synthesis['headline'][:80]}...")

    # 4. Build PPTX
    print("[4] Generating PowerPoint presentation...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover_slide(prs, agg, synthesis)
    print("    Slide 1: Cover - done")
    build_executive_summary_slide(prs, agg, synthesis)
    print("    Slide 2: Executive Summary - done")
    build_trends_slide(prs, agg, synthesis)
    print("    Slide 3: Trends - done")
    build_risks_slide(prs, agg, synthesis)
    print("    Slide 4: Emerging Risks - done")
    build_deepdive_slide(prs, agg)
    print("    Slide 5: Project Deep-Dives - done")
    build_recommendations_slide(prs, synthesis)
    print("    Slide 6: Recommendations - done")
    build_appendix_slide(prs, agg)
    print("    Slide 7: Appendix - done")

    # 5. Save
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    prs.save(PPTX_OUTPUT)
    print(f"\n[5] Presentation saved: {PPTX_OUTPUT}")
    print("=" * 60)
    print("  DONE! Open monthly_presentation.pptx in PowerPoint.")
    print("=" * 60)


if __name__ == "__main__":
    main()
